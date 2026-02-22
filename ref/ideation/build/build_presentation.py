#!/usr/bin/env python3
"""
Build script for Team of Teams Agent Skill presentation.
Generates a PowerPoint (.pptx) from the ideation session vision document.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BG       = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT_BLUE   = RGBColor(0x00, 0x96, 0xD6)   # bright blue
ACCENT_TEAL   = RGBColor(0x00, 0xB4, 0xA0)   # teal
ACCENT_AMBER  = RGBColor(0xFF, 0xA7, 0x26)   # amber/gold
ACCENT_RED    = RGBColor(0xE8, 0x4D, 0x4D)   # soft red
ACCENT_PURPLE = RGBColor(0x9B, 0x72, 0xCF)   # purple
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY      = RGBColor(0x99, 0x99, 0x99)
VERY_DARK     = RGBColor(0x12, 0x12, 0x22)
CARD_BG       = RGBColor(0x22, 0x22, 0x3A)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.05
    return shape


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_paragraph(tf, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_accent_line(slide, left, top, width, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bullet_list(tf, items, font_size=13, color=WHITE, bold=False):
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = f"\u2022  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold


def card_with_header(slide, left, top, width, height, header_text, body_lines, header_color=ACCENT_BLUE, body_font=12):
    """Create a card with a coloured header strip and body text."""
    card = add_shape_rect(slide, left, top, width, height, CARD_BG, border_color=header_color)
    # Header bar
    add_shape_rect(slide, left, top, width, Pt(28), header_color)
    # Header text
    tb = add_textbox(slide, left + Inches(0.15), top + Pt(2), width - Inches(0.3), Pt(26))
    set_text(tb.text_frame, header_text, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    # Body
    tb2 = add_textbox(slide, left + Inches(0.15), top + Pt(34), width - Inches(0.3), height - Pt(40))
    tf = tb2.text_frame
    tf.word_wrap = True
    add_bullet_list(tf, body_lines, font_size=body_font, color=LIGHT_GRAY)
    return card


# ============================================================================
# SLIDE BUILDERS
# ============================================================================

def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Decorative top accent
    add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Pt(6), ACCENT_BLUE)

    # Title
    tb = add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2))
    set_text(tb.text_frame, "Team of Teams Agent Skill", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    tb2 = add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8))
    set_text(tb2.text_frame, "A Company OS for Organizational Decision-Making", font_size=24, color=ACCENT_TEAL, bold=False, alignment=PP_ALIGN.CENTER)

    # Separator line
    add_accent_line(slide, Inches(5), Inches(4.0), Inches(3.333), ACCENT_AMBER)

    # Meta info
    tb3 = add_textbox(slide, Inches(1), Inches(4.4), Inches(11), Inches(1.5))
    tf = tb3.text_frame
    set_text(tf, "Multi-Agent Ideation Session", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "2026-02-22", font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(12))
    add_paragraph(tf, "Claude Code  |  Agent Teams", font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(8))


def build_overview_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Header
    tb = add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
    set_text(tb.text_frame, "The Vision", font_size=32, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(0.95), Inches(2), ACCENT_BLUE)

    # Core thesis card
    card = add_shape_rect(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(1.6), CARD_BG, border_color=ACCENT_BLUE)
    tb2 = add_textbox(slide, Inches(0.9), Inches(1.4), Inches(11.5), Inches(1.5))
    tf = tb2.text_frame
    tf.word_wrap = True
    set_text(tf, "Core Thesis", font_size=16, color=ACCENT_BLUE, bold=True)
    add_paragraph(tf,
        "A complete organizational reasoning engine -- a boardroom in a box -- that gives any founder, leader, "
        "or team the analytical power of a full executive committee. It transforms Claude Code into a Company OS "
        "where disagreement is engineered, not accidental, and every decision surfaces the tensions that matter.",
        font_size=14, color=LIGHT_GRAY, space_before=Pt(8))

    # Governing principle
    tb3 = add_textbox(slide, Inches(0.6), Inches(3.15), Inches(12), Inches(0.5))
    tf3 = tb3.text_frame
    set_text(tf3, "Governing Principle:", font_size=14, color=ACCENT_AMBER, bold=True)
    add_paragraph(tf3, "Structural dissent produces richer analysis than consensus ever could.", font_size=14, color=WHITE, space_before=Pt(2))

    # Three pillars overview
    pillar_data = [
        ("PILLAR 1: STRUCTURE", "Cascading Deliberation Engine", ACCENT_BLUE,
         ["Five-phase decision cascade", "CEO frames, C-suite decomposes, team leads analyze", "Hybrid agent architecture: 8 agents + 29 subagents", "Engineered dissent: 4 skeptics, 2 advocates, 1 systemic, 1 synthesizer"]),
        ("PILLAR 2: EXPERIENCE", "Engagement Model", ACCENT_TEAL,
         ["Three tiers: Hallway Question, Working Session, Board Meeting", "Five CEO decision modes (Guardian, Pioneer, Architect, Analyst, Sentinel)", "Multi-mode comparison: 1 analysis, N syntheses", "SMB-first: bias toward lightweight engagement"]),
        ("PILLAR 3: QUALITY", "Cognitive Forcing & Prompt Architecture", ACCENT_PURPLE,
         ["Three-layer prompt system prevents voice collapse", "29 unique analytical frameworks with mandatory output templates", "Forcing questions: pre-mortem, adversarial empathy, devil's advocate", "Heterogeneous model tiering: Haiku / Sonnet / Opus"]),
    ]

    x_positions = [Inches(0.6), Inches(4.75), Inches(8.9)]
    card_w = Inches(3.85)

    for i, (label, title, color, bullets) in enumerate(pillar_data):
        x = x_positions[i]
        y = Inches(3.85)
        h = Inches(3.35)
        card_with_header(slide, x, y, card_w, h, label, bullets, header_color=color, body_font=11)
        # Pillar title inside card
        ttb = add_textbox(slide, x + Inches(0.15), y + Pt(30), card_w - Inches(0.3), Pt(22))
        set_text(ttb.text_frame, title, font_size=13, color=color, bold=True)


def build_pillar1_slide_a(prs):
    """Pillar 1 part A: Five-phase cascade and hybrid architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    tb = add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
    set_text(tb.text_frame, "Pillar 1: The Cascading Deliberation Engine", font_size=28, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(0.9), Inches(3), ACCENT_BLUE)

    # Five phases -- left column
    phases = [
        ("Phase 0", "Shared Consciousness Broadcast", "CEO broadcasts issue context to all C-suite simultaneously -- shared picture before independent reasoning."),
        ("Phase 1", "CEO Frames & Routes", "Decompose into evaluation dimensions, classify decision type, route to relevant C-suite with explicit exclusion reasoning."),
        ("Phase 2", "C-Suite Dispatches Down", "Each exec translates CEO framing into domain-specific sub-questions -- analytical decomposition, not forwarding."),
        ("Phase 3", "Team Leads Produce Findings", "Narrow, specialist analysis through unique analytical frameworks and mandatory output templates."),
        ("Phase 4", "C-Suite Synthesizes Up", "Domain recommendations with confidence levels. Internal contradictions flagged as signals, not averaged away."),
        ("Phase 4.5", "Pre-Mortem Challenge", "Each C-suite reviews ALL peer recommendations: 'Assume failure in 12 months -- what caused it?' (Tier 3 only)"),
        ("Phase 5", "CEO Deliberation", "Map all recommendations, identify fault lines, apply Decision Mode, produce final Decision Record."),
    ]

    y_start = Inches(1.2)
    phase_h = Pt(42)
    gap = Pt(6)
    left_col_w = Inches(6.8)

    for i, (phase_num, phase_title, desc) in enumerate(phases):
        y = y_start + i * (phase_h + gap)
        # Phase number badge
        badge = add_shape_rect(slide, Inches(0.6), y, Inches(0.9), phase_h, ACCENT_BLUE)
        badge.text_frame.clear()
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = badge.text_frame.paragraphs[0].add_run()
        run.text = phase_num
        run.font.size = Pt(10)
        run.font.color.rgb = WHITE
        run.font.bold = True
        badge.text_frame.paragraphs[0].space_before = Pt(4)

        # Title + desc
        tb = add_textbox(slide, Inches(1.65), y, left_col_w - Inches(1.05), phase_h)
        tf = tb.text_frame
        tf.word_wrap = True
        set_text(tf, phase_title, font_size=13, color=WHITE, bold=True)
        add_paragraph(tf, desc, font_size=10, color=LIGHT_GRAY, space_before=Pt(2), space_after=Pt(0))

    # Right column: hybrid architecture
    right_x = Inches(7.6)
    right_w = Inches(5.2)

    tb_h = add_textbox(slide, right_x, Inches(1.2), right_w, Inches(0.4))
    set_text(tb_h.text_frame, "Hybrid Agent Architecture", font_size=18, color=ACCENT_TEAL, bold=True)
    add_accent_line(slide, right_x, Inches(1.65), Inches(2), ACCENT_TEAL)

    # Tier 1 card
    card_with_header(slide, right_x, Inches(1.9), right_w, Inches(1.5), "TIER 1: AGENT TEAM MEMBERS", [
        "CEO + 7 C-suite officers (COO, CFO, CTO, CISO, CAO, VP Sales, VP Delivery)",
        "Real teammates spawned by skill orchestrator",
        "Direct messaging between agents",
        "Model: Sonnet (C-suite), Opus (CEO)"
    ], header_color=ACCENT_TEAL, body_font=11)

    # Tier 2 card
    card_with_header(slide, right_x, Inches(3.65), right_w, Inches(1.5), "TIER 2: CUSTOM SUBAGENTS", [
        "29 team lead specialists in .claude/agents/team-leads/",
        "Each C-suite agent invokes their team leads",
        "Own context window, cost-efficient model (Haiku)",
        "Restricted tool access, structured result reporting"
    ], header_color=ACCENT_AMBER, body_font=11)

    # Benefits
    tb_ben = add_textbox(slide, right_x, Inches(5.4), right_w, Inches(1.8))
    tf_ben = tb_ben.text_frame
    tf_ben.word_wrap = True
    set_text(tf_ben, "Why This Matters", font_size=14, color=ACCENT_BLUE, bold=True)
    benefits = [
        "Real isolation: each perspective gets its own context window",
        "Cost optimization: Haiku for analysis, Sonnet/Opus for synthesis",
        "True management simulation: dispatch, collect, synthesize"
    ]
    add_bullet_list(tf_ben, benefits, font_size=11, color=LIGHT_GRAY)


def build_pillar1_slide_b(prs):
    """Pillar 1 part B: Engineered dissent and org roster."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    tb = add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
    set_text(tb.text_frame, "Engineered Dissent & Organizational Roster", font_size=28, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(0.9), Inches(3), ACCENT_BLUE)

    # Dissent model -- left half
    roles = [
        ("CEO", "Synthesizer", "Frame, listen, weigh, and decide", ACCENT_BLUE),
        ("COO", "Skeptic", "Can we actually do this with what we have?", ACCENT_RED),
        ("CFO", "Skeptic", "Find the costs not in the proposal", ACCENT_RED),
        ("CTO", "Advocate", "What does this make possible?", ACCENT_TEAL),
        ("CISO", "Skeptic", "Change introduces risk -- org's immune system", ACCENT_RED),
        ("VP Sales", "Advocate", "How does this help us sell more, faster?", ACCENT_TEAL),
        ("VP Delivery", "Skeptic", "What do we sacrifice from existing commitments?", ACCENT_RED),
        ("CAO", "Systemic", "Can the org absorb this?", ACCENT_PURPLE),
    ]

    y = Inches(1.3)
    row_h = Pt(36)
    for role_name, disposition, mandate, color in roles:
        # Role badge
        badge = add_shape_rect(slide, Inches(0.6), y, Inches(1.1), row_h, color)
        badge.text_frame.clear()
        p = badge.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = role_name
        r.font.size = Pt(11)
        r.font.color.rgb = WHITE
        r.font.bold = True
        p.space_before = Pt(5)

        # Disposition
        tb_d = add_textbox(slide, Inches(1.85), y, Inches(1.1), row_h)
        set_text(tb_d.text_frame, disposition, font_size=11, color=color, bold=True)

        # Mandate
        tb_m = add_textbox(slide, Inches(3.0), y, Inches(3.6), row_h)
        set_text(tb_m.text_frame, mandate, font_size=11, color=LIGHT_GRAY)

        y += row_h + Pt(4)

    # Balance summary
    tb_bal = add_textbox(slide, Inches(0.6), y + Pt(4), Inches(6), Inches(0.4))
    set_text(tb_bal.text_frame, "Balance: 4 skeptics  |  2 advocates  |  1 systemic  |  1 synthesizer", font_size=12, color=ACCENT_AMBER, bold=True)

    # Right side: org roster
    right_x = Inches(7.2)
    right_w = Inches(5.6)

    tb_t = add_textbox(slide, right_x, Inches(1.1), right_w, Inches(0.4))
    set_text(tb_t.text_frame, "Full Roster: 29 Team Leads", font_size=18, color=ACCENT_TEAL, bold=True)

    roster = [
        ("COO", "Operations Mgr, Process/Quality, Vendor/Procurement, Facilities", ACCENT_RED),
        ("CFO", "Controller, FP&A, Treasury/Cash, AP/AR, Tax Lead", ACCENT_RED),
        ("CTO", "Engineering, Infrastructure/DevOps, Data/Analytics, Product/UX", ACCENT_TEAL),
        ("CISO", "Security Ops, Compliance/GRC, Identity & Access, Security Architecture", ACCENT_RED),
        ("VP Sales", "Sales Ops, Account Mgmt, Business Dev, Sales Enablement", ACCENT_TEAL),
        ("VP Delivery", "Project/Program Mgr, Resource Mgr, Client Success, QA/Delivery", ACCENT_RED),
        ("CAO", "HR/People Ops, Legal/Contracts, Admin/Policy, Corp Comms", ACCENT_PURPLE),
    ]

    y2 = Inches(1.6)
    for exec_title, leads, color in roster:
        card_h = Inches(0.7)
        add_shape_rect(slide, right_x, y2, right_w, card_h, CARD_BG, border_color=color)
        # Exec label
        badge2 = add_shape_rect(slide, right_x + Inches(0.1), y2 + Pt(6), Inches(1.0), Pt(22), color)
        badge2.text_frame.clear()
        p2 = badge2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = exec_title
        r2.font.size = Pt(9)
        r2.font.color.rgb = WHITE
        r2.font.bold = True
        # Leads text
        tb_l = add_textbox(slide, right_x + Inches(1.25), y2 + Pt(4), right_w - Inches(1.5), Pt(30))
        set_text(tb_l.text_frame, leads, font_size=10, color=LIGHT_GRAY)

        y2 += card_h + Pt(4)

    # Decision type routing table at bottom right
    tb_rt = add_textbox(slide, right_x, y2 + Pt(8), right_w, Inches(0.3))
    set_text(tb_rt.text_frame, "Decision-Type Routing", font_size=14, color=ACCENT_AMBER, bold=True)

    routing = [
        "Strategic: CEO, CFO, CTO, VP Sales",
        "Operational: CEO, COO, VP Delivery",
        "Financial: CEO, CFO, COO",
        "Technical: CEO, CTO, CISO",
        "Personnel: CEO, CAO, COO, VP Delivery",
        "Compliance/Risk: CEO, CISO, CAO, CFO",
    ]
    tb_r = add_textbox(slide, right_x, y2 + Pt(30), right_w, Inches(1.2))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    add_bullet_list(tf_r, routing, font_size=10, color=LIGHT_GRAY)


def build_pillar2_slide_a(prs):
    """Pillar 2: Three tiers and decision modes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    tb = add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
    set_text(tb.text_frame, "Pillar 2: The Engagement Model", font_size=28, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(0.9), Inches(3), ACCENT_TEAL)

    # Three tiers
    tiers = [
        ("TIER 1", "Hallway Question", "/consult [role]: [question]",
         "Direct consult with one C-suite agent. No CEO, no routing, no team leads. Quick, opinionated, domain-specific.",
         "Advisory Note (3-5 sentences)", ACCENT_TEAL),
        ("TIER 2", "Working Session", "/panel [roles]: [issue]",
         "CEO routes to 2-4 C-suite members. Each does domain analysis with team lead perspectives. Lightweight synthesis.",
         "Panel Assessment (~1 page)", ACCENT_BLUE),
        ("TIER 3", "Board Meeting", "/deliberate: [issue]",
         "Full five-phase cascade. All relevant C-suite activated. Full team lead analysis. Full CEO deliberation.",
         "Complete Decision Record (3-5 pages)", ACCENT_PURPLE),
    ]

    x_pos = [Inches(0.6), Inches(4.75), Inches(8.9)]
    card_w = Inches(3.85)
    tier_y = Inches(1.2)
    tier_h = Inches(2.8)

    for i, (tier_label, name, invocation, desc, output, color) in enumerate(tiers):
        x = x_pos[i]
        add_shape_rect(slide, x, tier_y, card_w, tier_h, CARD_BG, border_color=color)
        # Header bar
        add_shape_rect(slide, x, tier_y, card_w, Pt(32), color)
        tb_h = add_textbox(slide, x + Inches(0.12), tier_y + Pt(3), card_w - Inches(0.24), Pt(28))
        tf_h = tb_h.text_frame
        set_text(tf_h, f"{tier_label}: {name}", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Invocation
        tb_inv = add_textbox(slide, x + Inches(0.15), tier_y + Pt(38), card_w - Inches(0.3), Pt(22))
        set_text(tb_inv.text_frame, invocation, font_size=11, color=ACCENT_AMBER, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        tb_desc = add_textbox(slide, x + Inches(0.15), tier_y + Pt(62), card_w - Inches(0.3), Inches(1.2))
        tf_desc = tb_desc.text_frame
        tf_desc.word_wrap = True
        set_text(tf_desc, desc, font_size=11, color=LIGHT_GRAY)

        # Output
        tb_out = add_textbox(slide, x + Inches(0.15), tier_y + Pt(120), card_w - Inches(0.3), Pt(24))
        set_text(tb_out.text_frame, f"Output: {output}", font_size=10, color=color, bold=True)

    # SMB-first note
    tb_smb = add_textbox(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(0.5))
    tf_smb = tb_smb.text_frame
    tf_smb.word_wrap = True
    set_text(tf_smb, "SMB-First Default:", font_size=13, color=ACCENT_AMBER, bold=True)
    r = tf_smb.paragraphs[0].add_run()
    r.text = "  Bias toward Tier 1. Most SMB decisions are fast and informal -- the skill should match that tempo."
    r.font.size = Pt(13)
    r.font.color.rgb = LIGHT_GRAY

    # Decision modes -- bottom section
    tb_dm = add_textbox(slide, Inches(0.6), Inches(4.85), Inches(12), Inches(0.4))
    set_text(tb_dm.text_frame, "Five CEO Decision Modes", font_size=18, color=ACCENT_TEAL, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(5.25), Inches(2), ACCENT_TEAL)

    modes = [
        ("Guardian", "Risk-averse (MaxiMin)", "Weights skeptics heavily. High bar for approval.", ACCENT_RED),
        ("Pioneer", "Growth-oriented (MaxiMax)", "Weights advocates heavily. Skeptic concerns = engineering problems.", ACCENT_TEAL),
        ("Architect", "Consensus (Behavioral)", "Weights fault lines. Seeks widest organizational support.", ACCENT_BLUE),
        ("Analyst", "Data-driven (Hurwicz)", "Weights confidence levels. 'Defer' is a legitimate outcome. DEFAULT.", ACCENT_AMBER),
        ("Sentinel", "Regret-minimizing (MiniMax Regret)", "Weights strongest objection. Favors survivable paths.", ACCENT_PURPLE),
    ]

    mode_w = Inches(2.38)
    mode_h = Inches(1.8)
    mode_x_positions = [Inches(0.6) + i * (mode_w + Inches(0.15)) for i in range(5)]
    mode_y = Inches(5.45)

    for i, (name, theory, desc, color) in enumerate(modes):
        x = mode_x_positions[i]
        add_shape_rect(slide, x, mode_y, mode_w, mode_h, CARD_BG, border_color=color)
        # Name
        tb_n = add_textbox(slide, x + Inches(0.1), mode_y + Pt(6), mode_w - Inches(0.2), Pt(22))
        set_text(tb_n.text_frame, name, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        # Theory
        tb_th = add_textbox(slide, x + Inches(0.1), mode_y + Pt(30), mode_w - Inches(0.2), Pt(20))
        set_text(tb_th.text_frame, theory, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
        # Desc
        tb_d = add_textbox(slide, x + Inches(0.1), mode_y + Pt(54), mode_w - Inches(0.2), Inches(0.9))
        tf_d = tb_d.text_frame
        tf_d.word_wrap = True
        set_text(tf_d, desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def build_pillar2_slide_b(prs):
    """Pillar 2 part B: Multi-mode comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    tb = add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
    set_text(tb.text_frame, "Multi-Mode Comparison & Comparative Decision Record", font_size=28, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(0.9), Inches(3), ACCENT_TEAL)

    # Key insight
    card = add_shape_rect(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.2), CARD_BG, border_color=ACCENT_AMBER)
    tb_k = add_textbox(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.0))
    tf_k = tb_k.text_frame
    tf_k.word_wrap = True
    set_text(tf_k, "Key Efficiency Insight", font_size=16, color=ACCENT_AMBER, bold=True)
    add_paragraph(tf_k,
        "Domain analysis is mode-independent. Multi-mode comparison runs domain analysis once (the expensive part) "
        "and CEO synthesis multiple times (cheap, single-agent passes). Cost: approximately 1.1x a single deliberation "
        "for 5x the strategic insight.",
        font_size=13, color=LIGHT_GRAY, space_before=Pt(8))

    # Invocation examples
    tb_inv = add_textbox(slide, Inches(0.6), Inches(2.7), Inches(5.5), Inches(0.4))
    set_text(tb_inv.text_frame, "Invocation Patterns", font_size=18, color=ACCENT_TEAL, bold=True)

    invocations = [
        "/deliberate guardian: should we acquire CompetitorX?",
        "/deliberate guardian vs pioneer: should we acquire CompetitorX?",
        "/deliberate all-modes: should we acquire CompetitorX?",
        "/consult cfo guardian: can we afford this?",
        "/panel pioneer finance tech: should we build this feature?",
    ]
    tb_inv2 = add_textbox(slide, Inches(0.6), Inches(3.15), Inches(5.5), Inches(2.5))
    tf_inv = tb_inv2.text_frame
    tf_inv.word_wrap = True
    for i, inv in enumerate(invocations):
        if i == 0:
            p = tf_inv.paragraphs[0]
        else:
            p = tf_inv.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = inv
        run.font.size = Pt(12)
        run.font.color.rgb = ACCENT_AMBER
        run.font.bold = True

    # Comparative Decision Record structure -- right side
    right_x = Inches(6.6)
    right_w = Inches(6.1)

    tb_cdr = add_textbox(slide, right_x, Inches(2.7), right_w, Inches(0.4))
    set_text(tb_cdr.text_frame, "Comparative Decision Record", font_size=18, color=ACCENT_BLUE, bold=True)

    sections = [
        "Executive Summary -- one paragraph per mode",
        "Shared Analysis -- domain analyses presented once",
        "Mode Comparisons -- per-mode synthesis with key factors",
        "Divergence Analysis -- where modes agree/diverge",
        "The Key Choice -- the values question underneath",
    ]
    card_with_header(slide, right_x, Inches(3.15), right_w, Inches(2.5),
                     "RECORD STRUCTURE", sections, header_color=ACCENT_BLUE, body_font=12)

    # Mode sensitivity
    tb_ms = add_textbox(slide, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.3))
    tf_ms = tb_ms.text_frame
    tf_ms.word_wrap = True
    set_text(tf_ms, "Mode Sensitivity -- A Novel Signal", font_size=16, color=ACCENT_PURPLE, bold=True)
    add_paragraph(tf_ms,
        "If all modes produce the same decision, the right answer does not depend on the user's risk appetite -- "
        "the evidence speaks for itself. If modes diverge dramatically, the user knows their personal risk posture "
        "is the deciding factor, not the analysis. This transforms a binary recommendation into a nuanced view of "
        "the decision landscape.",
        font_size=13, color=LIGHT_GRAY, space_before=Pt(8))


def build_pillar3_slide(prs):
    """Pillar 3: Cognitive forcing and prompt architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    tb = add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
    set_text(tb.text_frame, "Pillar 3: Cognitive Forcing & Prompt Architecture", font_size=28, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(0.9), Inches(3), ACCENT_PURPLE)

    # Voice collapse problem
    tb_vc = add_textbox(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(1.0))
    tf_vc = tb_vc.text_frame
    tf_vc.word_wrap = True
    set_text(tf_vc, "The Voice Collapse Problem", font_size=16, color=ACCENT_RED, bold=True)
    add_paragraph(tf_vc,
        "Telling one LLM to 'think as 5 different people' produces 5 paragraphs of the same reasoning "
        "with different vocabulary. Simple role assignments show zero measurable improvement. "
        "The solution: force each perspective to use different analytical methods producing structurally different outputs.",
        font_size=12, color=LIGHT_GRAY, space_before=Pt(6))

    # Three-layer architecture -- three cards
    layers = [
        ("LAYER 1", "Team Lead Subagent Definitions",
         "29 analytical framework packages in .claude/agents/team-leads/",
         ["Unique analytical framework per role",
          "Mandatory output template (structurally different per role)",
          "Three forcing questions: pre-mortem, adversarial empathy, devil's advocate",
          "Accountability framing and blind spot declarations",
          "Subagent config: model, tools, max turns, permission mode"],
         ACCENT_TEAL),
        ("LAYER 2", "C-Suite Agent Prompts",
         "7 domain orchestration templates",
         ["Mode A: Tier 1 direct consult (quick, opinionated)",
          "Mode B: Tier 2/3 full analysis (dispatch subagents, synthesize)",
          "Mode C: Phase 4.5 challenge (cross-peer pre-mortem)",
          "Role identity, mandate, team composition, natural tension partners",
          "Domain-level forcing questions and synthesis instructions"],
         ACCENT_BLUE),
        ("LAYER 3", "CEO Synthesis",
         "Modular decision mode system",
         ["Fixed: fault line mapping, determinative perspective identification",
          "Swappable: Decision Mode module adjusts synthesis weighting",
          "Same analysis, different weighting = different decisions",
          "Cross-domain pre-mortem integration",
          "Randomized recommendation ordering (anti-anchoring)"],
         ACCENT_PURPLE),
    ]

    layer_y = Inches(2.5)
    layer_h = Inches(2.8)
    layer_w = Inches(3.85)
    x_positions = [Inches(0.6), Inches(4.75), Inches(8.9)]

    for i, (label, title, subtitle, bullets, color) in enumerate(layers):
        x = x_positions[i]
        add_shape_rect(slide, x, layer_y, layer_w, layer_h, CARD_BG, border_color=color)
        # Header
        add_shape_rect(slide, x, layer_y, layer_w, Pt(32), color)
        tb_lh = add_textbox(slide, x + Inches(0.1), layer_y + Pt(4), layer_w - Inches(0.2), Pt(26))
        set_text(tb_lh.text_frame, label, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        # Title
        tb_lt = add_textbox(slide, x + Inches(0.12), layer_y + Pt(36), layer_w - Inches(0.24), Pt(22))
        set_text(tb_lt.text_frame, title, font_size=12, color=color, bold=True)
        # Subtitle
        tb_ls = add_textbox(slide, x + Inches(0.12), layer_y + Pt(56), layer_w - Inches(0.24), Pt(18))
        set_text(tb_ls.text_frame, subtitle, font_size=10, color=MID_GRAY)
        # Bullets
        tb_lb = add_textbox(slide, x + Inches(0.12), layer_y + Pt(76), layer_w - Inches(0.24), Inches(1.6))
        tf_lb = tb_lb.text_frame
        tf_lb.word_wrap = True
        add_bullet_list(tf_lb, bullets, font_size=10, color=LIGHT_GRAY)

    # Bias susceptibility table at bottom
    tb_bias = add_textbox(slide, Inches(0.6), Inches(5.55), Inches(12), Inches(0.3))
    set_text(tb_bias.text_frame, "Bias & Role Susceptibility Mitigations", font_size=16, color=ACCENT_AMBER, bold=True)

    biases = [
        ("Skeptics (CISO, CFO, COO, VP Delivery)", "Sycophancy softens objections", "\"Your value is in surfacing concerns, not being agreeable.\""),
        ("Advocates (CTO, VP Sales)", "Under-weight genuine constraints", "Must name strongest objection to their own position."),
        ("Synthesizer (CEO)", "False balance / anchoring bias", "Explicit weight rationale. Randomize recommendation order."),
        ("Systemic (CAO)", "Unfalsifiable vagueness", "Require concrete indicators: specific policies, teams, precedents."),
    ]

    bias_y = Inches(5.9)
    col_w = Inches(3.025)
    for i, (role, risk, mitigation) in enumerate(biases):
        x = Inches(0.6) + i * (col_w + Inches(0.1))
        tb_b = add_textbox(slide, x, bias_y, col_w, Inches(1.3))
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        set_text(tf_b, role, font_size=10, color=ACCENT_AMBER, bold=True)
        add_paragraph(tf_b, f"Risk: {risk}", font_size=9, color=ACCENT_RED, space_before=Pt(3))
        add_paragraph(tf_b, mitigation, font_size=9, color=LIGHT_GRAY, space_before=Pt(2))


def build_implementation_slide(prs):
    """Implementation specification: file structure, subagent format, model tiering."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    tb = add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
    set_text(tb.text_frame, "Implementation Specification", font_size=28, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(0.9), Inches(3), ACCENT_AMBER)

    # File structure -- left column
    tb_fs = add_textbox(slide, Inches(0.6), Inches(1.2), Inches(4.5), Inches(0.4))
    set_text(tb_fs.text_frame, "File Structure", font_size=18, color=ACCENT_BLUE, bold=True)

    fs_lines = [
        ".claude/agents/ceo.md",
        ".claude/agents/c-suite/ (7 agent definitions)",
        ".claude/agents/team-leads/ (7 domain dirs, 29 subagents)",
        ".claude/skills/team-of-teams/skill.md",
        ".claude/skills/team-of-teams/templates/",
        "  decision-record.md, advisory-note.md",
        "  panel-assessment.md, decision-space-map.md",
        ".claude/skills/team-of-teams/config/",
        "  routing-table.md, company-profile.md",
        "  decision-modes.md",
    ]

    tb_fsl = add_textbox(slide, Inches(0.6), Inches(1.65), Inches(4.5), Inches(3.5))
    tf_fsl = tb_fsl.text_frame
    tf_fsl.word_wrap = True
    for i, line in enumerate(fs_lines):
        if i == 0:
            p = tf_fsl.paragraphs[0]
        else:
            p = tf_fsl.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(11)
        run.font.color.rgb = ACCENT_TEAL if not line.startswith(" ") else LIGHT_GRAY
        run.font.bold = not line.startswith(" ")

    # Subagent format -- middle column
    right_x = Inches(5.5)
    right_w = Inches(7.2)

    tb_sf = add_textbox(slide, right_x, Inches(1.2), right_w, Inches(0.4))
    set_text(tb_sf.text_frame, "Subagent Definition Format", font_size=18, color=ACCENT_TEAL, bold=True)

    format_items = [
        "YAML front matter: name, description, model (haiku), tools, maxTurns",
        "Role Identity: title, C-suite parent, domain scope",
        "Analytical Framework: specific methodology (e.g., GAAP Compliance Assessment)",
        "Mandatory Output Template: structured fields unique to this role",
        "Three Forcing Questions: pre-mortem, adversarial empathy, devil's advocate",
        "Blind Spot Declaration: explicit scope exclusions",
        "Accountability Framing: analysis reviewed by C-suite parent alongside peers",
    ]

    card_with_header(slide, right_x, Inches(1.65), right_w, Inches(2.8),
                     "EACH OF 29 TEAM LEAD FILES CONTAINS:", format_items, header_color=ACCENT_TEAL, body_font=11)

    # Model tiering -- bottom section
    tb_mt = add_textbox(slide, Inches(0.6), Inches(4.8), Inches(12), Inches(0.4))
    set_text(tb_mt.text_frame, "Heterogeneous Model Tiering", font_size=18, color=ACCENT_PURPLE, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(5.2), Inches(2), ACCENT_PURPLE)

    tiers_data = [
        ("Team Lead Subagents", "Haiku", "Narrow, focused analysis. Cost-efficient. Model diversity improves variety.", ACCENT_TEAL),
        ("C-Suite Agents", "Sonnet", "Domain decomposition, delegation, and synthesis. Moderate complexity.", ACCENT_BLUE),
        ("CEO Agent", "Opus", "Cross-domain synthesis, fault-line analysis, decision-making. Highest quality.", ACCENT_PURPLE),
    ]

    tier_w = Inches(3.85)
    tier_x = [Inches(0.6), Inches(4.75), Inches(8.9)]
    tier_y = Inches(5.4)

    for i, (layer, model, rationale, color) in enumerate(tiers_data):
        x = tier_x[i]
        add_shape_rect(slide, x, tier_y, tier_w, Inches(1.3), CARD_BG, border_color=color)
        tb_l = add_textbox(slide, x + Inches(0.15), tier_y + Pt(6), tier_w - Inches(0.3), Pt(22))
        set_text(tb_l.text_frame, layer, font_size=13, color=color, bold=True)
        tb_m = add_textbox(slide, x + Inches(0.15), tier_y + Pt(28), tier_w - Inches(0.3), Pt(24))
        set_text(tb_m.text_frame, f"Model: {model}", font_size=16, color=WHITE, bold=True)
        tb_r = add_textbox(slide, x + Inches(0.15), tier_y + Pt(54), tier_w - Inches(0.3), Inches(0.6))
        tf_r = tb_r.text_frame
        tf_r.word_wrap = True
        set_text(tf_r, rationale, font_size=10, color=LIGHT_GRAY)

    # Company profile config note
    tb_cp = add_textbox(slide, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.4))
    tf_cp = tb_cp.text_frame
    tf_cp.word_wrap = True
    set_text(tf_cp, "Company Profile:", font_size=12, color=ACCENT_AMBER, bold=True)
    r_cp = tf_cp.paragraphs[0].add_run()
    r_cp.text = "  Configurable roster, default decision mode, tier escalation thresholds, and industry-specific frameworks. Default: mid-market IT services, 200-500 employees."
    r_cp.font.size = Pt(12)
    r_cp.font.color.rgb = LIGHT_GRAY


def build_future_slide(prs):
    """Future directions and boundaries."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    tb = add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
    set_text(tb.text_frame, "Boundaries & Future Directions", font_size=28, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(0.9), Inches(3), ACCENT_AMBER)

    # Current scope boundaries
    tb_scope = add_textbox(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.3))
    set_text(tb_scope.text_frame, "Current Scope (v1)", font_size=18, color=ACCENT_BLUE, bold=True)

    scope_items = [
        "Single-issue deliberation (one question at a time)",
        "Point-in-time analysis (no persistent memory across sessions)",
        "One-round pre-mortem (no multi-round debate)",
        "Fixed analytical frameworks (no community contribution yet)",
        "Default company profile: mid-market IT services, 200-500 employees",
    ]

    card_with_header(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.4),
                     "IN SCOPE", scope_items, header_color=ACCENT_BLUE, body_font=12)

    # Future directions
    future_items = [
        ("Decision Revisitation", "Revisit decisions when assumptions change. Re-run only affected domain analyses. Produce Decision Amendments rather than full new records. Transforms from point-in-time tool to decision lifecycle manager.", ACCENT_TEAL),
        ("Scenario Forking Engine", "Parallel cascades with different parameters producing comparative decision landscapes. Not just 'what should we do?' but 'what happens if X vs. Y vs. Z?'", ACCENT_PURPLE),
        ("Organizational Stress Test", "Invert the cascade: bottom-up threat and opportunity discovery. Each C-suite generates issues from their domain. Output: Organizational Health Report.", ACCENT_AMBER),
        ("Institutional Memory", "Persistent knowledge across sessions. Previous Decision Records inform current analysis. The system learns company patterns.", ACCENT_BLUE),
        ("Cross-Functional Deliberation", "Full Phase 4.5 where C-suite agents debate each other's recommendations -- structured around specific fault lines to prevent problem drift.", ACCENT_RED),
        ("Framework Versioning", "Analytical framework packages versioned and improved independently. Industry-specific variants. Community contributions. Skill becomes a platform.", ACCENT_TEAL),
    ]

    right_x = Inches(6.8)
    right_w = Inches(5.9)
    tb_ft = add_textbox(slide, right_x, Inches(1.3), right_w, Inches(0.3))
    set_text(tb_ft.text_frame, "Future Directions (v2+)", font_size=18, color=ACCENT_TEAL, bold=True)

    y = Inches(1.7)
    card_h = Inches(0.88)
    for title, desc, color in future_items:
        add_shape_rect(slide, right_x, y, right_w, card_h, CARD_BG, border_color=color)
        tb_f = add_textbox(slide, right_x + Inches(0.15), y + Pt(4), right_w - Inches(0.3), card_h - Pt(8))
        tf_f = tb_f.text_frame
        tf_f.word_wrap = True
        set_text(tf_f, title, font_size=12, color=color, bold=True)
        add_paragraph(tf_f, desc, font_size=10, color=LIGHT_GRAY, space_before=Pt(3))
        y += card_h + Pt(6)

    # Onboarding stress test note at bottom left
    tb_ob = add_textbox(slide, Inches(0.6), Inches(4.4), Inches(5.8), Inches(0.3))
    set_text(tb_ob.text_frame, "Onboarding Stress Test", font_size=16, color=ACCENT_PURPLE, bold=True)

    onboard_items = [
        "Run a representative cross-functional issue through full Tier 3 cascade on first configuration",
        "Validates agent configuration for the company type",
        "Calibrates user expectations on depth and quality",
        "Surfaces profile issues (e.g., irrelevant roles for company type)",
    ]
    card_with_header(slide, Inches(0.6), Inches(4.75), Inches(5.8), Inches(2.0),
                     "INITIALIZATION PROTOCOL", onboard_items, header_color=ACCENT_PURPLE, body_font=11)


def build_closing_slide(prs):
    """Closing: cross-cutting themes, open questions, next steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    tb = add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6))
    set_text(tb.text_frame, "Cross-Cutting Themes & Next Steps", font_size=28, color=WHITE, bold=True)
    add_accent_line(slide, Inches(0.6), Inches(0.9), Inches(3), ACCENT_BLUE)

    # Cross-cutting themes
    themes = [
        ("Engineered Dissent > Consensus",
         "Structural tensions between skeptics, advocates, and systemic perspectives are features, not bugs. "
         "The Fault Line Analysis is the single most valuable artifact the system produces."),
        ("Analytical Method Diversity Prevents Voice Collapse",
         "Each of 29 team leads uses a distinct analytical framework with a mandatory output template. "
         "Different methods produce structurally different outputs that cannot collapse into one voice."),
        ("Heterogeneous Architecture Mirrors Real Organizations",
         "Three model tiers (Haiku/Sonnet/Opus) map to three organizational layers. "
         "Model diversity improves analytical variety beyond cost savings."),
        ("Decision Modes Separate Analysis from Values",
         "The same domain analysis, weighted differently, produces different decisions. "
         "Mode Sensitivity reveals whether the evidence or the user's posture drives the outcome."),
        ("SMB-First Design Philosophy",
         "Bias toward lightweight engagement. Tier 1 is the daily habit; Tier 3 is the deliberate escalation. "
         "The skill matches the tempo of how real SMB leaders make decisions."),
    ]

    y = Inches(1.2)
    theme_h = Inches(0.85)
    theme_w = Inches(12.1)
    colors = [ACCENT_BLUE, ACCENT_TEAL, ACCENT_PURPLE, ACCENT_AMBER, ACCENT_RED]

    for i, (title, desc) in enumerate(themes):
        color = colors[i % len(colors)]
        add_shape_rect(slide, Inches(0.6), y, theme_w, theme_h, CARD_BG, border_color=color)
        # Color strip on left
        add_shape_rect(slide, Inches(0.6), y, Pt(4), theme_h, color)
        tb_th = add_textbox(slide, Inches(0.85), y + Pt(4), theme_w - Inches(0.5), theme_h - Pt(8))
        tf_th = tb_th.text_frame
        tf_th.word_wrap = True
        set_text(tf_th, title, font_size=13, color=color, bold=True)
        add_paragraph(tf_th, desc, font_size=11, color=LIGHT_GRAY, space_before=Pt(3))
        y += theme_h + Pt(6)

    # Next steps
    y += Pt(8)
    tb_ns = add_textbox(slide, Inches(0.6), y, Inches(12), Inches(0.3))
    set_text(tb_ns.text_frame, "Suggested Next Steps", font_size=18, color=ACCENT_TEAL, bold=True)

    steps = [
        "1. Build skill scaffolding and CEO agent definition with five-phase cascade protocol",
        "2. Implement CFO domain as reference implementation (5 complete team lead packages)",
        "3. Extend to remaining 6 C-suite agents and 24 team lead subagents",
        "4. Implement invocation commands (/consult, /panel, /deliberate, /evaluate) with tier routing",
        "5. Add Decision Mode overlays (Guardian, Pioneer, Architect, Analyst, Sentinel) and multi-mode comparison",
        "6. Run onboarding stress test to validate all 29 perspectives produce genuinely different analysis",
    ]

    tb_steps = add_textbox(slide, Inches(0.6), y + Inches(0.35), Inches(12.1), Inches(1.8))
    tf_steps = tb_steps.text_frame
    tf_steps.word_wrap = True
    add_bullet_list(tf_steps, steps, font_size=12, color=LIGHT_GRAY)

    # Footer
    tb_footer = add_textbox(slide, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3))
    set_text(tb_footer.text_frame, "Team of Teams Agent Skill  |  Multi-Agent Ideation Session  |  2026-02-22",
             font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_title_slide(prs)          # Slide 1: Title
    build_overview_slide(prs)       # Slide 2: Overview / core thesis
    build_pillar1_slide_a(prs)      # Slide 3: Cascading Deliberation Engine
    build_pillar1_slide_b(prs)      # Slide 4: Engineered Dissent & Roster
    build_pillar2_slide_a(prs)      # Slide 5: Engagement Model (tiers + modes)
    build_pillar2_slide_b(prs)      # Slide 6: Multi-Mode Comparison
    build_pillar3_slide(prs)        # Slide 7: Cognitive Forcing & Prompt Arch
    build_implementation_slide(prs) # Slide 8: Implementation Spec
    build_future_slide(prs)         # Slide 9: Boundaries & Future
    build_closing_slide(prs)        # Slide 10: Cross-cutting Themes & Next Steps

    out_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    out_path = os.path.join(out_dir, "PRESENTATION_team-of-teams.pptx")
    prs.save(out_path)
    print(f"Presentation saved to: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
